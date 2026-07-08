from __future__ import annotations

import io
import os
import re
import uuid

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ..config import settings
from .artifacts import resolve_image_bytes
from .base import ToolContext, ToolDef, ToolResult

# Subdirectory (under UPLOAD_DIR) where generated downloadable files are written.
GENERATED_SUBDIR = "generated"

FONT_HEAD = "Segoe UI Semibold"
FONT_BODY = "Segoe UI"

# 16:9 widescreen canvas.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


class Theme:
    def __init__(
        self,
        bg: str,
        accent: str,
        title_text: str,
        subtitle_text: str,
        heading: str,
        body: str,
        bullet: str,
    ) -> None:
        self.bg = RGBColor.from_string(bg)
        self.accent = RGBColor.from_string(accent)
        self.title_text = RGBColor.from_string(title_text)
        self.subtitle_text = RGBColor.from_string(subtitle_text)
        self.heading = RGBColor.from_string(heading)
        self.body = RGBColor.from_string(body)
        self.bullet = RGBColor.from_string(bullet)


THEMES: dict[str, Theme] = {
    "indigo": Theme("1E1B4B", "6366F1", "FFFFFF", "C7D2FE", "1E1B4B", "334155", "6366F1"),
    "midnight": Theme("0F172A", "38BDF8", "FFFFFF", "94A3B8", "0F172A", "334155", "38BDF8"),
    "teal": Theme("064E3B", "10B981", "FFFFFF", "A7F3D0", "064E3B", "334155", "10B981"),
    "plum": Theme("4A1D48", "D946EF", "FFFFFF", "F5D0FE", "4A1D48", "3F3F46", "D946EF"),
    "slate": Theme("1F2937", "F59E0B", "FFFFFF", "D1D5DB", "1F2937", "374151", "F59E0B"),
    "crimson": Theme("450A0A", "EF4444", "FFFFFF", "FECACA", "450A0A", "3F3F46", "EF4444"),
}
DEFAULT_THEME = "indigo"

# Palette used to color chart series / slices.
PALETTE = [
    "6366F1", "10B981", "F59E0B", "EF4444", "3B82F6",
    "8B5CF6", "EC4899", "14B8A6", "F97316", "84CC16",
]


def _generated_dir() -> str:
    path = os.path.join(settings.UPLOAD_DIR, GENERATED_SUBDIR)
    os.makedirs(path, exist_ok=True)
    return path


def _safe_download_name(title: str) -> str:
    base = re.sub(r"[^A-Za-z0-9 _-]+", "", title or "").strip().replace(" ", "-")
    base = re.sub(r"-{2,}", "-", base).strip("-")
    if not base:
        base = "presentation"
    return f"{base[:60]}.pptx"


def _no_border(shape) -> None:
    shape.line.fill.background()
    shape.shadow.inherit = False


def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    _no_border(shape)


class PptxGenerateTool:
    definition = ToolDef(
        name="generate_pptx",
        description=(
            "Create a downloadable, professionally STYLED Microsoft PowerPoint (.pptx) "
            "presentation from a structured outline that you author. Use this whenever "
            "the user asks for a PowerPoint, slide deck, or presentation. Provide a "
            "title and a list of slides (each with a slide title and bullet points). "
            "Optionally pick a color theme. The tool applies a themed title slide, "
            "accent bars, styled fonts, and footers automatically, and returns a "
            "Markdown download link — include that link verbatim in your reply. "
            "IMPORTANT: When the user asks for any visualization, chart, graph, pie "
            "chart, or bar chart, you MUST render a REAL chart by populating that "
            "slide's 'chart' field with a type plus numeric 'categories' and 'series' "
            "data. Do NOT merely write 'Bar chart' or 'Pie chart' as a slide title or "
            "bullet — always provide the actual data so an editable chart is drawn. "
            "Supply concrete numbers (research or estimate them if needed)."
        ),
        parameters={
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    "description": "Title of the presentation (shown on the title slide).",
                },
                "subtitle": {
                    "type": "string",
                    "description": "Optional subtitle for the title slide.",
                },
                "theme": {
                    "type": "string",
                    "enum": list(THEMES.keys()),
                    "description": "Color theme for the deck. Defaults to 'indigo'.",
                },
                "slides": {
                    "type": "array",
                    "description": "The content slides, in order.",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {
                                "type": "string",
                                "description": "The slide's heading.",
                            },
                            "bullets": {
                                "type": "array",
                                "description": "Bullet points for the slide body.",
                                "items": {"type": "string"},
                            },
                            "notes": {
                                "type": "string",
                                "description": "Optional speaker notes for the slide.",
                            },
                            "image": {
                                "type": "string",
                                "description": (
                                    "Optional image to place on the slide: a data: URI, "
                                    "an http(s) image URL, or a generated /api/files/<name> "
                                    "link (e.g. from generate_image)."
                                ),
                            },
                            "chart": {
                                "type": "object",
                                "description": (
                                    "Optional native chart to render on the slide. "
                                    "Provide real data so an editable chart is drawn."
                                ),
                                "properties": {
                                    "type": {
                                        "type": "string",
                                        "enum": ["pie", "bar", "column", "line", "doughnut"],
                                        "description": "Chart type.",
                                    },
                                    "title": {
                                        "type": "string",
                                        "description": "Optional chart title.",
                                    },
                                    "categories": {
                                        "type": "array",
                                        "items": {"type": "string"},
                                        "description": "Category labels (x-axis / pie slices).",
                                    },
                                    "series": {
                                        "type": "array",
                                        "description": (
                                            "One or more data series. Pie/doughnut use "
                                            "the first series only."
                                        ),
                                        "items": {
                                            "type": "object",
                                            "properties": {
                                                "name": {"type": "string"},
                                                "values": {
                                                    "type": "array",
                                                    "items": {"type": "number"},
                                                },
                                            },
                                            "required": ["values"],
                                        },
                                    },
                                },
                                "required": ["type", "categories", "series"],
                            },
                        },
                        "required": ["title"],
                    },
                },
            },
            "required": ["title", "slides"],
        },
    )

    # -- slide builders -------------------------------------------------------

    def _title_slide(self, prs, theme: Theme, title: str, subtitle: str) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # Full-bleed background.
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        _fill(bg, theme.bg)
        # Accent bar.
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(3.05), Inches(1.6), Inches(0.12)
        )
        _fill(bar, theme.accent)
        # Title.
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(3.3), Inches(11.5), Inches(2.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = FONT_HEAD
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = theme.title_text
        # Subtitle.
        if subtitle:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = subtitle
            r2.font.name = FONT_BODY
            r2.font.size = Pt(20)
            r2.font.color.rgb = theme.subtitle_text

    def _content_slide(
        self, prs, theme: Theme, index: int, deck_title: str, item: dict
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # White canvas.
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        _fill(bg, RGBColor.from_string("FFFFFF"))
        # Left accent rail.
        rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), SLIDE_H)
        _fill(rail, theme.accent)

        # Heading + underline.
        head = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(1.0))
        htf = head.text_frame
        htf.word_wrap = True
        hp = htf.paragraphs[0]
        hr = hp.add_run()
        hr.text = (item.get("title") or "").strip() or "Slide"
        hr.font.name = FONT_HEAD
        hr.font.size = Pt(30)
        hr.font.bold = True
        hr.font.color.rgb = theme.heading
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.45), Inches(2.2), Inches(0.06)
        )
        _fill(underline, theme.accent)

        # Chart (optional) + bullets layout.
        chart_spec = item.get("chart") if isinstance(item.get("chart"), dict) else None
        image_ref = (item.get("image") or "").strip() if isinstance(item.get("image"), str) else ""
        side_visual = bool(chart_spec) or bool(image_ref)
        bullets = item.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        bullets = [str(b).strip() for b in bullets if str(b).strip()]

        chart_left = chart_top = chart_w = chart_h = None
        if side_visual and bullets:
            body_w = Inches(5.6)
            chart_left, chart_top, chart_w, chart_h = (
                Inches(6.7), Inches(1.9), Inches(6.0), Inches(4.8),
            )
        elif side_visual:
            body_w = Inches(11.8)
            chart_left, chart_top, chart_w, chart_h = (
                Inches(1.2), Inches(1.9), Inches(11.0), Inches(4.9),
            )
        else:
            body_w = Inches(11.8)

        body = slide.shapes.add_textbox(Inches(0.75), Inches(1.85), body_w, Inches(4.9))
        btf = body.text_frame
        btf.word_wrap = True
        first = True
        for bullet in bullets:
            text = str(bullet).strip()
            if not text:
                continue
            p = btf.paragraphs[0] if first else btf.add_paragraph()
            p.space_after = Pt(12)
            dot = p.add_run()
            dot.text = "\u25AA  "
            dot.font.name = FONT_BODY
            dot.font.size = Pt(18)
            dot.font.color.rgb = theme.bullet
            r = p.add_run()
            r.text = text
            r.font.name = FONT_BODY
            r.font.size = Pt(18)
            r.font.color.rgb = theme.body
            first = False

        # Chart (optional).
        if chart_spec and chart_left is not None:
            self._add_chart(
                slide, theme, chart_spec, chart_left, chart_top, chart_w, chart_h
            )
        elif image_ref and chart_left is not None:
            self._add_image(
                slide, image_ref, chart_left, chart_top, chart_w, chart_h
            )

        # Footer: deck title (left) + page number (right).
        foot = slide.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(9), Inches(0.4))
        fp = foot.text_frame.paragraphs[0]
        fr = fp.add_run()
        fr.text = deck_title
        fr.font.name = FONT_BODY
        fr.font.size = Pt(10)
        fr.font.color.rgb = RGBColor.from_string("94A3B8")

        num = slide.shapes.add_textbox(Inches(12.0), Inches(6.95), Inches(1.0), Inches(0.4))
        ntf = num.text_frame
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = ntf.paragraphs[0]
        np.alignment = PP_ALIGN.RIGHT
        nr = np.add_run()
        nr.text = str(index)
        nr.font.name = FONT_BODY
        nr.font.size = Pt(10)
        nr.font.bold = True
        nr.font.color.rgb = theme.accent

        notes = (item.get("notes") or "").strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _add_image(self, slide, ref: str, left, top, width, height) -> bool:
        """Embed an image (data URI / http URL / /api/files link) into a slide,
        centered within the given box and preserving aspect ratio."""
        try:
            data = resolve_image_bytes(ref)
            if not data:
                return False
            from PIL import Image as _PILImage

            iw, ih = _PILImage.open(io.BytesIO(data)).size
            box_w, box_h = int(width), int(height)
            scale = min(box_w / iw, box_h / ih)
            w = int(iw * scale)
            h = int(ih * scale)
            cx = int(left) + (box_w - w) // 2
            cy = int(top) + (box_h - h) // 2
            slide.shapes.add_picture(io.BytesIO(data), Emu(cx), Emu(cy), Emu(w), Emu(h))
            return True
        except Exception:  # noqa: BLE001
            return False

    def _add_chart(self, slide, theme: Theme, spec: dict, left, top, width, height) -> bool:
        """Render a native (editable) chart from a chart spec. Returns True on success."""
        try:
            ctype = (spec.get("type") or "column").strip().lower()
            cats = [str(c) for c in (spec.get("categories") or [])]
            raw_series = spec.get("series") or []
            norm: list[tuple[str, list[float]]] = []
            for s in raw_series:
                if not isinstance(s, dict):
                    continue
                vals: list[float] = []
                for v in s.get("values") or []:
                    try:
                        vals.append(float(v))
                    except (TypeError, ValueError):
                        vals.append(0.0)
                norm.append((str(s.get("name") or "Series"), vals))
            if not cats or not norm:
                return False

            chart_type = {
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "pie": XL_CHART_TYPE.PIE,
                "doughnut": XL_CHART_TYPE.DOUGHNUT,
                "line": XL_CHART_TYPE.LINE_MARKERS,
            }.get(ctype, XL_CHART_TYPE.COLUMN_CLUSTERED)
            is_pie = chart_type in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT)

            data = CategoryChartData()
            data.categories = cats
            if is_pie:
                data.add_series(norm[0][0], norm[0][1])
            else:
                for name, vals in norm:
                    data.add_series(name, vals)

            frame = slide.shapes.add_chart(chart_type, left, top, width, height, data)
            chart = frame.chart
            chart.font.size = Pt(11)
            chart.font.name = FONT_BODY

            title = (spec.get("title") or "").strip()
            chart.has_title = bool(title)
            if title:
                chart.chart_title.text_frame.text = title

            plot = chart.plots[0]
            if is_pie:
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.RIGHT
                chart.legend.include_in_layout = False
                plot.has_data_labels = True
                dl = plot.data_labels
                dl.show_percentage = True
                dl.show_value = False
                dl.number_format = "0%"
                dl.number_format_is_linked = False
                for i, point in enumerate(plot.series[0].points):
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor.from_string(
                        PALETTE[i % len(PALETTE)]
                    )
            else:
                multi = len(norm) > 1
                chart.has_legend = multi
                if multi:
                    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                    chart.legend.include_in_layout = False
                    for i, s in enumerate(plot.series):
                        s.format.fill.solid()
                        s.format.fill.fore_color.rgb = RGBColor.from_string(
                            PALETTE[i % len(PALETTE)]
                        )
                else:
                    plot.vary_by_categories = True
                    for i, point in enumerate(plot.series[0].points):
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = RGBColor.from_string(
                            PALETTE[i % len(PALETTE)]
                        )
            return True
        except Exception:  # noqa: BLE001
            return False

    # -- entrypoint -----------------------------------------------------------

    async def run(self, args: dict, ctx: ToolContext) -> ToolResult:
        title = (args.get("title") or "").strip() or "Presentation"
        subtitle = (args.get("subtitle") or "").strip()
        slides = args.get("slides") or []
        theme = THEMES.get((args.get("theme") or "").strip().lower(), THEMES[DEFAULT_THEME])
        if not isinstance(slides, list) or not slides:
            return ToolResult(
                content="No slides were provided, so no presentation was created.",
                citations=None,
            )

        try:
            prs = Presentation()
            prs.slide_width = SLIDE_W
            prs.slide_height = SLIDE_H

            self._title_slide(prs, theme, title, subtitle)
            n = 0
            for item in slides:
                if not isinstance(item, dict):
                    continue
                n += 1
                self._content_slide(prs, theme, n, title, item)

            file_id = uuid.uuid4().hex
            stored_name = f"{file_id}.pptx"
            path = os.path.join(_generated_dir(), stored_name)
            prs.save(path)
        except Exception as exc:  # noqa: BLE001
            return ToolResult(content=f"Failed to build presentation: {exc}", citations=None)

        download_name = _safe_download_name(title)
        url = f"/api/files/{stored_name}?name={download_name}"
        return ToolResult(
            content=(
                f"Created a {n}-slide PowerPoint presentation. "
                f"[\U0001F4E5 Download {download_name}]({url})"
            ),
            citations=None,
        )
